from pptx import Presentation
from pptx.util import Inches, Pt

# Your text
text = """Leader: We thank God for fathers, who pointed us to God, and who have loved and cared for us.
People: We give thanks and praise for fathers young and old.
--
Leader: We pray for young fathers, newly embracing their vocation;
People: may they find courage and perseverance to balance work, family and faith in joy and sacrifice.
--
(rest of your text)"""

# Split by --
sections = [section.strip() for section in text.split('--') if section.strip()]

# Load existing presentation or create new
prs = Presentation('Presentation1.pptx')  # or Presentation() for new

# Add slides
for section in sections:
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add text box
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    text_frame = textbox.text_frame
    text_frame.text = section
    
    # Format text
    text_frame.paragraphs[0].font.size = Pt(32)
    text_frame.paragraphs[0].font.name = 'Calibri'

prs.save('Updated_Presentation.pptx')
