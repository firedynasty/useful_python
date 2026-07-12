#!/usr/bin/env python3
"""
Generate prayer/creed slides from a text file.

Beige background, black text, centered. Lines are grouped into slides
using blank lines or '---' lines as section breaks. Within each section,
lines are paired (2 per slide). Single remaining lines get their own slide.

Usage:
  python make_prayer_slides.py apostles_creed.txt
  python make_prayer_slides.py lords_prayer.txt
  python make_prayer_slides.py -o output.pptx lords_prayer.txt
"""

import argparse
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def parse_into_slides(filepath, lines_per_slide=2):
    """Parse text file into slide groups.

    Blank lines and '---' lines act as hard slide breaks.
    Within each section, lines are grouped into chunks of `lines_per_slide`.
    """
    with open(filepath, 'r') as f:
        raw = f.readlines()

    slides = []
    current = []

    def flush():
        if current:
            for i in range(0, len(current), lines_per_slide):
                slides.append(current[i:i + lines_per_slide])

    for line in raw:
        stripped = line.strip()
        if not stripped or stripped.startswith('---'):
            flush()
            current = []
        else:
            current.append(stripped)

    flush()
    return slides


def create_slides(text_file, title=None, output_path=None):
    """Create a PPTX presentation from a prayer/creed text file."""
    slides_data = parse_into_slides(text_file)

    if title is None:
        title = os.path.splitext(os.path.basename(text_file))[0].replace('_', ' ').title()

    if not slides_data:
        print(f"  No text found in {text_file}, skipping.")
        return

    if output_path is None:
        base = os.path.splitext(os.path.basename(text_file))[0]
        output_path = os.path.join(os.path.dirname(text_file) or '.', f'{base}.pptx')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    FONT_SIZE = Pt(44)
    BLACK = RGBColor(0, 0, 0)
    BEIGE = RGBColor(235, 224, 199)
    FONT_NAME = 'Arial'

    def add_beige_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BEIGE

    def add_slide_number(slide, number):
        box_size = Inches(0.55)
        margin = Inches(0.2)
        tb = slide.shapes.add_textbox(
            slide_width - box_size - margin,
            slide_height - box_size - margin,
            box_size, box_size
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(18)
        p.font.name = FONT_NAME
        p.font.color.rgb = RGBColor(120, 110, 95)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # --- Title slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_beige_background(slide)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), slide_width - Inches(2), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.name = FONT_NAME
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # --- Content slides ---
    slide_num = 2
    for group in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_beige_background(slide)

        left = Inches(0.5)
        width = slide_width - Inches(1)
        top = Inches(2.0)
        height = Inches(5.0)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.clear()

        for i, line in enumerate(group):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = FONT_SIZE
            p.font.name = FONT_NAME
            p.font.color.rgb = BLACK
            p.space_after = Pt(8)

        add_slide_number(slide, slide_num)
        slide_num += 1

    # --- Blank end slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_beige_background(slide)

    prs.save(output_path)
    print(f"  Created: {output_path} ({len(prs.slides)} slides)")


def main():
    parser = argparse.ArgumentParser(description='Generate prayer/creed slides')
    parser.add_argument('file', help='Prayer .txt file')
    parser.add_argument('-t', '--title', default=None, help='Slide title (default: from filename)')
    parser.add_argument('-o', '--output', default=None, help='Output .pptx path')
    args = parser.parse_args()

    if not os.path.exists(args.file):
        print(f"Error: file not found: {args.file}")
        sys.exit(1)

    print(f"Processing: {args.file}")
    create_slides(args.file, args.title, args.output)


if __name__ == '__main__':
    main()
