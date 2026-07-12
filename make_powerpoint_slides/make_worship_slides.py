#!/usr/bin/env python3
"""
Generate worship lyric slides from OpenSong-format text files.

Each slide shows all lines of the current section. The pair currently
being sung is underlined+bold; the rest are regular weight. The section
label (e.g. [Chorus]) appears in the top-right corner, and a slide
number box appears in the bottom-right.

Also generates a companion printable PDF lyric sheet with slide numbers.

Usage:
  python make_worship_slides.py worship/create_in_me_clean_heart.txt
  python make_worship_slides.py worship/*.txt
"""

import argparse
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib import colors
from reportlab.pdfbase import pdfmetrics


def parse_lyrics(filepath):
    """Parse a lyrics file top-to-bottom.

    - Lines starting with . are comments (skipped)
    - Lines like [Chorus], [Verse1] are section labels (tracked)
    - || is a page break within a section (starts a new pair)
    - Everything else is a lyric line

    Returns a list of sections, each:
        (label, display_lines, pairs)
    where pairs is a list of lists of bold line indices into display_lines.
    """
    with open(filepath, 'r') as f:
        raw_lines = f.read().split('\n')

    sections = []
    current_label = ''
    current_lines = []

    def flush_section():
        if not current_lines:
            return
        # Build pairs of 2 lines respecting || breaks
        pairs = []
        group = []
        for idx, line in enumerate(current_lines):
            if line == '||':
                for i in range(0, len(group), 2):
                    pairs.append(group[i:i+2])
                group = []
            else:
                group.append(idx)
        for i in range(0, len(group), 2):
            pairs.append(group[i:i+2])

        # Filter || from display lines, remap indices
        display_lines = [l for l in current_lines if l != '||']
        orig_to_display = {}
        di = 0
        for oi, l in enumerate(current_lines):
            if l != '||':
                orig_to_display[oi] = di
                di += 1
        remapped_pairs = [[orig_to_display[idx] for idx in pair] for pair in pairs]
        sections.append((current_label, display_lines, remapped_pairs))

    for line in raw_lines:
        stripped = line.strip()
        if stripped.startswith('.'):
            continue
        if stripped.startswith('[') and stripped.endswith(']'):
            flush_section()
            current_label = stripped
            current_lines = []
            continue
        if not stripped:
            continue
        if stripped == '||':
            current_lines.append('||')
            continue
        current_lines.append(stripped)

    flush_section()
    return sections


def create_slides(lyrics_file, background_path, output_path=None):
    """Create a PPTX presentation from a lyrics file."""
    sections = parse_lyrics(lyrics_file)
    title = os.path.splitext(os.path.basename(lyrics_file))[0].replace('_', ' ').title()

    if not sections:
        print(f"  No lyrics found in {lyrics_file}, skipping.")
        return

    if output_path is None:
        base = os.path.splitext(os.path.basename(lyrics_file))[0]
        output_path = os.path.join(os.path.dirname(lyrics_file), f'{base}.pptx')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    FONT_SIZE = Pt(44)
    WHITE = RGBColor(255, 255, 255)
    LABEL_COLOR = RGBColor(200, 200, 210)
    NUM_COLOR = RGBColor(180, 180, 195)
    FONT_NAME = 'Arial'

    def add_background(slide):
        slide.shapes.add_picture(background_path, Emu(0), Emu(0), slide_width, slide_height)

    def add_section_label(slide, label):
        if not label:
            return
        tb = slide.shapes.add_textbox(
            slide_width - Inches(3.5), Inches(0.3), Inches(3.2), Inches(0.6)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(22)
        p.font.name = FONT_NAME
        p.font.color.rgb = LABEL_COLOR
        p.font.italic = True
        p.alignment = PP_ALIGN.RIGHT

    def add_slide_number(slide, number):
        """Small slide number box in the bottom-right corner."""
        box_size = Inches(0.55)
        margin = Inches(0.2)
        tb = slide.shapes.add_textbox(
            slide_width - box_size - margin,
            slide_height - box_size - margin,
            box_size, box_size
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(18)
        p.font.name = FONT_NAME
        p.font.color.rgb = NUM_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def add_lyrics_block(slide, all_lines, bold_indices):
        left = Inches(0.5)
        width = slide_width - Inches(1)
        top = Inches(1.5)
        height = Inches(5.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.clear()

        bold_set = set(bold_indices)

        for i, line in enumerate(all_lines):
            if not line:
                continue
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = FONT_SIZE
            p.font.name = FONT_NAME
            p.font.color.rgb = WHITE
            p.font.bold = i in bold_set
            p.font.underline = i in bold_set
            p.space_after = Pt(6)

    # --- Title slide (slide 1, not numbered) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), slide_width - Inches(2), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.name = FONT_NAME
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # --- Lyrics slides (numbered from 2) ---
    slide_num = 2
    for label, all_lines, pairs in sections:
        for bold_indices in pairs:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_background(slide)
            add_section_label(slide, label)
            add_lyrics_block(slide, all_lines, bold_indices)
            add_slide_number(slide, slide_num)
            slide_num += 1

    # --- Blank end slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    prs.save(output_path)
    print(f"  Created: {output_path} ({len(prs.slides)} slides)")
    return sections, title, output_path


def create_print_sheet(sections, title, output_path):
    """Generate a printable PDF lyric sheet with slide numbers."""
    page_w, page_h = A4
    margin = 2 * cm
    usable_w = page_w - 2 * margin

    c = rl_canvas.Canvas(output_path, pagesize=A4)

    # Fonts
    TITLE_SIZE = 18
    SECTION_SIZE = 10
    LYRIC_SIZE = 12
    NUM_SIZE = 9
    LINE_H = 0.65 * cm
    BOX_SIZE = 0.55 * cm
    GAP_AFTER_SECTION = 0.3 * cm
    GAP_AFTER_PAIR = 0.5 * cm

    y = page_h - margin

    def new_page():
        nonlocal y
        c.showPage()
        y = page_h - margin

    def check_space(needed):
        if y - needed < margin:
            new_page()

    # Title
    c.setFont("Helvetica-Bold", TITLE_SIZE)
    c.drawCentredString(page_w / 2, y, title)
    y -= 1.0 * cm

    # Thin rule under title
    c.setStrokeColor(colors.HexColor('#888888'))
    c.setLineWidth(0.5)
    c.line(margin, y, page_w - margin, y)
    y -= 0.5 * cm

    # Numbering: title slide = 1, lyrics start at 2
    slide_num = 2

    for label, all_lines, pairs in sections:
        # Section label
        check_space(0.8 * cm)
        c.setFont("Helvetica-Oblique", SECTION_SIZE)
        c.setFillColor(colors.HexColor('#555555'))
        c.drawString(margin, y, label)
        y -= GAP_AFTER_SECTION

        for pair_idx, bold_indices in enumerate(pairs):
            bold_set = set(bold_indices)
            n_lines = len(all_lines)
            needed = n_lines * LINE_H + GAP_AFTER_PAIR + 0.2 * cm
            check_space(needed)

            # Number box (top-left of the pair block)
            box_x = margin
            box_y = y - BOX_SIZE
            c.setStrokeColor(colors.HexColor('#333333'))
            c.setFillColor(colors.HexColor('#f0f0f0'))
            c.setLineWidth(0.6)
            c.roundRect(box_x, box_y, BOX_SIZE, BOX_SIZE, 2, stroke=1, fill=1)
            c.setFillColor(colors.HexColor('#222222'))
            c.setFont("Helvetica-Bold", NUM_SIZE)
            c.drawCentredString(box_x + BOX_SIZE / 2, box_y + (BOX_SIZE - NUM_SIZE * 0.72) / 2, str(slide_num))

            # Lyric lines indented past the box
            text_x = margin + BOX_SIZE + 0.3 * cm
            text_w = usable_w - BOX_SIZE - 0.3 * cm

            for i, line in enumerate(all_lines):
                if i in bold_set:
                    c.setFont("Helvetica-Bold", LYRIC_SIZE)
                    c.setFillColor(colors.black)
                else:
                    c.setFont("Helvetica", LYRIC_SIZE)
                    c.setFillColor(colors.HexColor('#444444'))
                c.drawString(text_x, y - i * LINE_H, line)

            y -= n_lines * LINE_H + GAP_AFTER_PAIR
            slide_num += 1

        y -= 0.3 * cm  # extra gap between sections

    c.save()
    print(f"  Created: {output_path}")


def main():
    parser = argparse.ArgumentParser(description='Generate worship lyric slides + print sheet')
    parser.add_argument('files', nargs='+', help='Lyrics .txt file(s)')
    parser.add_argument('-b', '--background',
                        default=os.path.join(os.path.dirname(__file__), 'worship', 'background_lighter.jpg'),
                        help='Background image path')
    parser.add_argument('-o', '--output-dir', default=None,
                        help='Output directory (default: same as input file)')
    args = parser.parse_args()

    if not os.path.exists(args.background):
        print(f"Error: Background image not found: {args.background}")
        sys.exit(1)

    for filepath in args.files:
        if not os.path.exists(filepath):
            print(f"Skipping {filepath}: file not found")
            continue
        print(f"Processing: {filepath}")
        base = os.path.splitext(os.path.basename(filepath))[0]
        out_dir = args.output_dir or os.path.dirname(filepath)

        pptx_out = os.path.join(out_dir, f'{base}.pptx')
        pdf_out = os.path.join(out_dir, f'{base}_print.pdf')

        result = create_slides(filepath, args.background, pptx_out)
        if result:
            sections, title, _ = result
            create_print_sheet(sections, title, pdf_out)


if __name__ == '__main__':
    main()
