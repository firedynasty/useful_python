from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import argparse
import os
import re

parser = argparse.ArgumentParser(description='Generate PowerPoint slides from a markdown announcements file.')
parser.add_argument('-i', '--input', required=True, help='Input markdown file (e.g. Sunday_service_annoucements.md)')
parser.add_argument('-o', '--output', help='Output .pptx file (defaults to input name with .pptx extension)')
args = parser.parse_args()

INPUT_FILE = args.input
OUTPUT_FILE = args.output or os.path.splitext(os.path.basename(INPUT_FILE))[0] + '.pptx'

TITLE_COLOR = RGBColor(0, 51, 102)       # Dark navy
BODY_COLOR = RGBColor(51, 51, 51)        # Dark gray
CHILD_COLOR = RGBColor(80, 80, 80)       # Medium gray
TITLE_SIZE = Pt(36)
BODY_SIZE = Pt(36)
CHILD_SIZE = Pt(36)

# --- Parse the markdown into slides ---
# **Bold Title** starts a new slide (becomes the slide heading).
# Bullets (* or - ) under a title become body text on that slide.

with open(INPUT_FILE, 'r') as f:
    lines = f.readlines()

slides_data = []
current_slide = None

for line in lines:
    stripped = line.rstrip()
    if not stripped:
        continue

    # Bold title line: **Title** — starts a new slide
    bold_match = re.match(r'^\*\*(.+?)\*\*$', stripped)
    if bold_match:
        if current_slide is not None:
            slides_data.append(current_slide)
        current_slide = {
            'title': bold_match.group(1).strip(),
            'children': []
        }
    # Bullet point: * or - at start of line — adds body text to current slide
    elif re.match(r'^[\*\-]\s', stripped) and current_slide is not None:
        bullet_text = re.sub(r'^[\*\-]\s+', '', stripped).strip()
        current_slide['children'].append(bullet_text)

if current_slide is not None:
    slides_data.append(current_slide)

print(f"Parsed {len(slides_data)} announcements from '{INPUT_FILE}'")

# --- Build the presentation ---
prs = Presentation()
# Set widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Announcement slides — each has a bold title + bullet points
for i, s in enumerate(slides_data):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(6.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Slide title (from **Bold Title**)
    p = tf.paragraphs[0]
    p.text = s['title']
    p.font.size = TITLE_SIZE
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(16)

    # Bullet points
    for child in s['children']:
        p = tf.add_paragraph()
        p.text = f"\u2022  {child}"
        p.font.size = Pt(36)
        p.font.color.rgb = BODY_COLOR
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(12)

    print(f"Slide {i+1}: {s['title']}")

prs.save(OUTPUT_FILE)
print(f"\nSaved '{OUTPUT_FILE}' with {len(prs.slides)} slides (1 title + {len(slides_data)} announcements)")
