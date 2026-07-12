#!/usr/bin/env python3
"""Create a PowerPoint presentation with one image per slide, sized to fill."""

import sys
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

SUPPORTED_FORMATS = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp')

def make_pptx(image_dir, output_path):
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    files = sorted(
        f for f in os.listdir(image_dir)
        if f.lower().endswith(SUPPORTED_FORMATS)
    )

    if not files:
        print(f"No images found in {image_dir}")
        sys.exit(1)

    for fname in files:
        fpath = os.path.realpath(os.path.join(image_dir, fname))

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        with Image.open(fpath) as img:
            img_w, img_h = img.size

        # Scale to cover the entire slide (fill, no letterbox)
        scale_w = slide_width / img_w
        scale_h = slide_height / img_h
        scale = max(scale_w, scale_h)

        new_w = int(img_w * scale)
        new_h = int(img_h * scale)
        left = (slide_width - new_w) // 2
        top = (slide_height - new_h) // 2

        slide.shapes.add_picture(fpath, Emu(left), Emu(top), Emu(new_w), Emu(new_h))

    prs.save(output_path)
    print(f"Saved {len(files)} slides to {output_path}")

if __name__ == '__main__':
    if len(sys.argv) != 3:
        print(f"Usage: {sys.argv[0]} <image_directory> <output.pptx>")
        sys.exit(1)
    make_pptx(sys.argv[1], sys.argv[2])
