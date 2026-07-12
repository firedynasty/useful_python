#!/usr/bin/env python3
"""
Generate Apostles' Creed slides from a text file.

Each slide shows 2-3 lines of the creed on a beige background with black text.
Blank lines in the input start a new slide group.

Usage:
  python make_creed_slides.py ../apostles_creed.txt
  python make_creed_slides.py -o output.pptx ../apostles_creed.txt
"""

import argparse
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# Slide groupings: list of (start_line, end_line) 0-indexed into non-empty lines
# Based on the natural phrasing of the Apostles' Creed
SLIDE_GROUPS = [
    [0, 1, 2],       # I believe in God, the Father almighty, Creator of heaven and earth,
    [3, 4],           # and in Jesus Christ... conceived by the Holy Spirit,
    [5, 6],           # born of the Virgin Mary, suffered under Pontius Pilate,
    [7, 8],           # was crucified, died and was buried; he descended into hell;
    [9],              # on the third day he rose again from the dead;
    [10, 11],         # he ascended into heaven... right hand of God the Father almighty;
    [12],             # from there he will come to judge the living and the dead.
    [13, 14],         # I believe in the Holy Spirit, the holy catholic Church,
    [15, 16],         # the communion of saints, the forgiveness of sins,
    [17, 18],         # the resurrection of the body, and life everlasting.
    [19],             # Amen.
]


def parse_creed(filepath):
    """Parse creed text file, returning non-empty lines."""
    with open(filepath, 'r') as f:
        lines = [line.strip() for line in f.readlines()]
    return [line for line in lines if line]


def create_slides(lyrics_file, output_path=None):
    """Create a PPTX presentation from the creed text."""
    all_lines = parse_creed(lyrics_file)
    title = os.path.splitext(os.path.basename(lyrics_file))[0].replace('_', ' ').title()

    if not all_lines:
        print(f"  No text found in {lyrics_file}, skipping.")
        return

    if output_path is None:
        base = os.path.splitext(os.path.basename(lyrics_file))[0]
        output_path = os.path.join(os.path.dirname(lyrics_file) or '.', f'{base}.pptx')

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    FONT_SIZE = Pt(44)
    BLACK = RGBColor(0, 0, 0)
    BEIGE = RGBColor(235, 224, 199)
    FONT_NAME = 'Arial'

    def add_beige_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BEIGE

    def add_slide_number(slide, number):
        box_size = Inches(0.55)
        margin = Inches(0.2)
        tb = slide.shapes.add_textbox(
            slide_width - box_size - margin,
            slide_height - box_size - margin,
            box_size, box_size
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(18)
        p.font.name = FONT_NAME
        p.font.color.rgb = RGBColor(120, 110, 95)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # --- Title slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_beige_background(slide)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), slide_width - Inches(2), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Apostles' Creed"
    p.font.size = Pt(60)
    p.font.name = FONT_NAME
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # --- Content slides ---
    # Use predefined groupings if they fit, otherwise fall back to 2-line groups
    if len(all_lines) >= max(max(g) for g in SLIDE_GROUPS) + 1:
        groups = SLIDE_GROUPS
    else:
        groups = [list(range(i, min(i + 2, len(all_lines)))) for i in range(0, len(all_lines), 2)]

    slide_num = 2
    for group in groups:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_beige_background(slide)

        left = Inches(0.5)
        width = slide_width - Inches(1)
        top = Inches(2.0)
        height = Inches(5.0)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.clear()

        for i, line_idx in enumerate(group):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = all_lines[line_idx]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = FONT_SIZE
            p.font.name = FONT_NAME
            p.font.color.rgb = BLACK
            p.space_after = Pt(8)

        add_slide_number(slide, slide_num)
        slide_num += 1

    # --- Blank end slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_beige_background(slide)

    prs.save(output_path)
    print(f"  Created: {output_path} ({len(prs.slides)} slides)")


def main():
    parser = argparse.ArgumentParser(description='Generate Apostles\' Creed slides')
    parser.add_argument('file', help='Creed .txt file')
    parser.add_argument('-o', '--output', default=None, help='Output .pptx path')
    args = parser.parse_args()

    if not os.path.exists(args.file):
        print(f"Error: file not found: {args.file}")
        sys.exit(1)

    print(f"Processing: {args.file}")
    create_slides(args.file, args.output)


if __name__ == '__main__':
    main()
