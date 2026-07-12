from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Your complete liturgical text
text = """Leader: We thank God for fathers, who pointed us to God, and who have loved and cared for us.
People: We give thanks and praise for fathers young and old.
--
Leader: We pray for young fathers, newly embracing their vocation;
People: may they find courage and perseverance to balance work, family and faith in joy and sacrifice.
--
Leader: We pray for fathers around the world whose children are lost or suffering;
People: may they know that the God of compassion walks with them in their sorrow.
--
Leader: We pray for men who are not fathers
People: but still mentor and guide us with fatherly love and advice.
--
Leader: We remember fathers, grandfathers, and great grandfathers
People: who are no longer with us but who have nourished us with their love and live forever in our memory.
--
Leader: And for God's love for us, his children, that continues to be the model for all faithful fathers:
ALL: we give thanks to God, our Heavenly Father. Amen."""

# Split by -- and clean up
sections = [section.strip() for section in text.split('--') if section.strip()]

# Load your existing presentation
try:
    prs = Presentation('Presentation1.pptx')
    print(f"Loaded existing presentation with {len(prs.slides)} slides")
except:
    # Create new if file doesn't exist
    prs = Presentation()
    print("Created new presentation")

# Add slides for each section
for i, section in enumerate(sections):
    # Use blank layout
    slide_layout = prs.slide_layouts[6]  # 6 is typically blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add text box with proper dimensions for wrapping
    left = Inches(0.5)
    top = Inches(1)
    width = Inches(9)    # Wide enough for content
    height = Inches(6)   # Tall enough for wrapped text
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    # Enable word wrap
    text_frame.word_wrap = True
    text_frame.auto_size = None  # Disable auto-sizing
    
    # Clear default paragraph and add formatted content
    text_frame.clear()
    
    lines = section.split('\n')
    for j, line in enumerate(lines):
        if line.strip():  # Skip empty lines
            if j == 0:
                p = text_frame.paragraphs[0]  # Use existing first paragraph
            else:
                p = text_frame.add_paragraph()
            
            p.text = line.strip()
            p.font.size = Pt(32)  # Large font for projection
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.LEFT
            
            # Color coding
            if line.startswith('Leader:'):
                p.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue
                p.font.bold = True
            elif line.startswith('People:'):
                p.font.color.rgb = RGBColor(0, 100, 0)  # Dark green
                p.font.bold = True
            elif line.startswith('ALL:'):
                p.font.color.rgb = RGBColor(139, 0, 0)  # Dark red
                p.font.bold = True
            
            # Add some space between Leader and People lines
            if j > 0:
                p.space_before = Pt(12)
    
    print(f"Added slide {i+1}: {section[:50]}...")

# Save the presentation
output_filename = 'Updated_Presentation.pptx'
prs.save(output_filename)
print(f"\nSaved presentation as '{output_filename}'")
print(f"Total slides: {len(prs.slides)}")
