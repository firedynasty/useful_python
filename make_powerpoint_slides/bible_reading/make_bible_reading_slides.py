#!/usr/bin/env python3
"""
Generate Bible reading slides matching the LCHS format.

Layout:
  - Outer background: light cream
  - Inner content box: darker tan/beige
  - Reference (italic) top-left of inner box
  - Bold dark-brown verse text with superscript verse numbers
  - Dark teal divider line at bottom
  - "Lutheran Church of the Holy Spirit | San Francisco" footer

Text file format:
  Line 1: .Reference (e.g. .John 10:1-10)  — dot prefix = reference
  Line 2: Optional heading (plain text, used as slide title if present)
  Remaining: Verse text with embedded numbers (e.g. "2 But he who...")
  Blank lines separate slide groups.

Usage:
  python make_bible_reading_slides.py ../bible_reading.txt


text file:
  .John 14:15-31  
  Jesus Promises the Holy Spirit                                                
  15 "If you love me, you will keep my commandments...
                                                                                
  So when copying from BibleGateway:                                            
                                                                              
  1. First line: Add the reference with a dot prefix: .John 14:15-31            
  2. Second line: The section heading (e.g. "Jesus Promises the Holy Spirit")
  3. Rest: Paste the verse text as-is       


⏺ Now fix the txt file — keep it simple, just paste continuously with blank
  lines only between paragraph groups:

  




"""

import argparse
import os
import re
import sys
import textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


# Colors from the screenshot
OUTER_BG = RGBColor(225, 220, 210)     # light cream outer
INNER_BG = RGBColor(206, 197, 178)     # darker tan inner box
TEXT_COLOR = RGBColor(75, 55, 30)       # dark brown text
REF_COLOR = RGBColor(95, 80, 55)       # slightly lighter brown for reference
TEAL = RGBColor(0, 70, 70)             # dark teal for line and footer
FOOTER_TEXT = "Lutheran Church of the Holy Spirit | San Francisco"


def parse_bible_reading(filepath):
    """Parse bible reading text file.

    Returns (reference, heading, paragraphs) where paragraphs is a list
    of text blocks (split by blank lines).
    """
    with open(filepath, 'r') as f:
        lines = f.readlines()

    reference = ''
    heading = ''
    text_lines = []

    for line in lines:
        stripped = line.strip()
        if stripped.startswith('.'):
            reference = stripped[1:].strip()
            continue
        if not reference:
            continue
        if not heading and stripped and not any(c.isdigit() and i == 0 for i, c in enumerate(stripped)):
            # First non-reference, non-verse line is the heading
            # But only if it doesn't start with a chapter number followed by verse text
            if not re.match(r'^\d+\s+["\u201c]', stripped):
                heading = stripped
                continue
        text_lines.append(line)

    # Join and split by blank lines into paragraphs
    full_text = ''.join(text_lines)
    raw_paragraphs = [p.strip() for p in full_text.split('\n\n') if p.strip()]
    # Collapse single newlines within each paragraph to spaces
    paragraphs = [re.sub(r'\s*\n\s*', ' ', p) for p in raw_paragraphs]

    return reference, heading, paragraphs


def split_into_slide_chunks(paragraphs, max_chars=249):
    """Split text into slides by character count, breaking at word boundaries.

    Treats the passage as a continuous reading flow — like paginating a book.
    Each slide gets up to max_chars characters, with the cut landing at the
    last space before the limit so no word is split mid-way.

    At 40pt bold Arial on this slide layout, ~200 chars fills ~7 lines safely.
    Increase max_chars for fewer, fuller slides (may risk slight overflow).
    """
    text = ' '.join(p for p in paragraphs if p)
    chunks = []
    while len(text) > max_chars:
        cut = text.rfind(' ', 0, max_chars + 1)
        if cut <= 0:
            cut = max_chars  # fallback: hard cut if no space found
        chunks.append(text[:cut].rstrip())
        text = text[cut:].lstrip()
    if text:
        chunks.append(text)
    return chunks


def add_verse_text_with_superscript(text_frame, text, font_size, font_name, text_color):
    """Add verse text with superscript verse numbers.

    Verse numbers are detected as standalone numbers (e.g. "2 But", "10 The").
    """
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # Split text into segments: verse numbers vs regular text
    # Match verse numbers that appear at start or after punctuation/space
    parts = re.split(r'(\b(\d{1,3})\s)(?=[A-Z"\u201c])', text)

    # Simpler approach: find all verse number positions
    tokens = re.split(r'(?<![:\-\d])(\b\d{1,3}\b)(?=\s+[A-Z"\u201c\u2018])', text)

    i = 0
    while i < len(tokens):
        token = tokens[i]
        # Check if this token is a verse number
        if re.match(r'^\d{1,3}$', token) and i + 1 < len(tokens):
            # Superscript verse number
            run = p.add_run()
            run.text = token
            run.font.size = Pt(font_size.pt * 0.55)
            run.font.name = font_name
            run.font.color.rgb = text_color
            run.font.bold = True
            # Set superscript via XML
            rPr = run._r.get_or_add_rPr()
            rPr.set('baseline', '40000')  # positive = superscript
            # Add a thin space after
            run2 = p.add_run()
            run2.text = ''
            run2.font.size = font_size
        else:
            # Regular text
            if token:
                run = p.add_run()
                run.text = token
                run.font.size = font_size
                run.font.name = font_name
                run.font.color.rgb = text_color
                run.font.bold = True
        i += 1


def create_slides(text_file, output_path=None):
    """Create Bible reading PPTX slides."""
    reference, heading, paragraphs = parse_bible_reading(text_file)

    if not paragraphs:
        print(f"  No text found in {text_file}, skipping.")
        return

    if output_path is None:
        base = os.path.splitext(os.path.basename(text_file))[0]
        output_path = os.path.join(os.path.dirname(text_file) or '.', f'{base}.pptx')

    ref_display = reference + ' ESV' if reference and 'ESV' not in reference else reference

    chunks = split_into_slide_chunks(paragraphs)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw = prs.slide_width
    sh = prs.slide_height

    FONT_NAME = 'Arial'
    BODY_FONT_SIZE = Pt(40)

    # Inner box dimensions
    box_left = Inches(0.6)
    box_top = Inches(0.5)
    box_width = sw - Inches(1.2)
    box_height = Inches(5.8)

    from pptx.enum.shapes import MSO_SHAPE

    def add_slide_chrome(slide):
        """Add the common outer bg, inner box, teal line, and footer."""
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = OUTER_BG

        inner_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            box_left, box_top, box_width, box_height
        )
        inner_box.fill.solid()
        inner_box.fill.fore_color.rgb = INNER_BG
        inner_box.line.fill.background()
        inner_box.adjustments[0] = 0.02

        # Dark teal divider line
        line_top = box_top + box_height + Inches(0.25)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            box_left, line_top,
            box_width, Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = TEAL
        line.line.fill.background()

        # Footer text
        footer_tb = slide.shapes.add_textbox(
            box_left, line_top + Inches(0.15),
            box_width, Inches(0.4)
        )
        tf = footer_tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = FOOTER_TEXT
        p.font.size = Pt(16)
        p.font.name = FONT_NAME
        p.font.color.rgb = TEAL
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

    # --- Title slide (with inner box, teal line, and footer) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_chrome(slide)

    # "God's Word Comes To Us" — top left of inner box, italic
    gwc_tb = slide.shapes.add_textbox(
        box_left + Inches(0.4), box_top + Inches(0.3),
        box_width - Inches(0.8), Inches(1.0)
    )
    tf = gwc_tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "God's Word Comes To Us"
    p.font.size = Pt(48)
    p.font.name = FONT_NAME
    p.font.color.rgb = TEXT_COLOR
    p.font.italic = True
    p.alignment = PP_ALIGN.LEFT

    # "Bible Reading" — below, smaller, regular
    br_tb = slide.shapes.add_textbox(
        box_left + Inches(0.4), box_top + Inches(1.25),
        box_width - Inches(0.8), Inches(0.5)
    )
    tf = br_tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "Bible Reading"
    p.font.size = Pt(22)
    p.font.name = FONT_NAME
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = False
    p.alignment = PP_ALIGN.LEFT

    # Reference (e.g. "John 10:1-10") — bold, large, left-aligned, mid-box
    ref_only = reference  # without ESV
    ref_tb = slide.shapes.add_textbox(
        box_left + Inches(0.4), box_top + Inches(2.5),
        box_width - Inches(0.8), Inches(1.2)
    )
    tf = ref_tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = ref_only
    p.font.size = Pt(64)
    p.font.name = FONT_NAME
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # "English Standard Version" — italic, below reference
    esv_tb = slide.shapes.add_textbox(
        box_left + Inches(0.4), box_top + Inches(3.8),
        box_width - Inches(0.8), Inches(0.9)
    )
    tf = esv_tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "English Standard Version"
    p.font.size = Pt(44)
    p.font.name = FONT_NAME
    p.font.color.rgb = TEXT_COLOR
    p.font.italic = True
    p.alignment = PP_ALIGN.LEFT

    # --- Content slides ---
    def make_slide(chunk_text, show_heading=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_chrome(slide)

        # Reference label (italic, top-left of inner box)
        ref_tb = slide.shapes.add_textbox(
            box_left + Inches(0.4), box_top + Inches(0.3),
            Inches(6), Inches(0.5)
        )
        tf = ref_tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = ref_display
        p.font.size = Pt(24)
        p.font.name = FONT_NAME
        p.font.color.rgb = REF_COLOR
        p.font.italic = True
        p.alignment = PP_ALIGN.LEFT

        if show_heading and heading:
            # Subtitle heading below reference
            sub_tb = slide.shapes.add_textbox(
                box_left + Inches(0.4), box_top + Inches(0.85),
                box_width - Inches(0.8), Inches(0.55)
            )
            tf = sub_tb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = heading
            p.font.size = Pt(26)
            p.font.name = FONT_NAME
            p.font.color.rgb = REF_COLOR
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT
            verse_top = box_top + Inches(1.55)
        else:
            verse_top = box_top + Inches(1.0)

        # Verse text (bold, dark brown)
        text_tb = slide.shapes.add_textbox(
            box_left + Inches(0.4), verse_top,
            box_width - Inches(0.8), box_height - (verse_top - box_top) - Inches(0.2)
        )
        tf = text_tb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        add_verse_text_with_superscript(tf, chunk_text, BODY_FONT_SIZE, FONT_NAME, TEXT_COLOR)

    for i, chunk in enumerate(chunks):
        make_slide(chunk, show_heading=(i == 0))

    prs.save(output_path)
    print(f"  Created: {output_path} ({len(prs.slides)} slides)")


def main():
    parser = argparse.ArgumentParser(description='Generate Bible reading slides')
    parser.add_argument('file', help='Bible reading .txt file')
    parser.add_argument('-o', '--output', default=None, help='Output .pptx path')
    args = parser.parse_args()

    if not os.path.exists(args.file):
        print(f"Error: file not found: {args.file}")
        sys.exit(1)

    print(f"Processing: {args.file}")
    create_slides(args.file, args.output)


if __name__ == '__main__':
    main()
