from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import argparse
import os
import re

parser = argparse.ArgumentParser(description='Generate PowerPoint slides from a markdown announcements file.')
parser.add_argument('-i', '--input', required=True, help='Input markdown file (e.g. Sunday_service_annoucements.md)')
parser.add_argument('-o', '--output', help='Output .pptx file (defaults to input name with .pptx extension)')
args = parser.parse_args()

INPUT_FILE = args.input
OUTPUT_FILE = args.output or os.path.splitext(os.path.basename(INPUT_FILE))[0] + '.pptx'

TITLE_COLOR = RGBColor(0, 51, 102)       # Dark navy
BODY_COLOR = RGBColor(51, 51, 51)        # Dark gray
CHILD_COLOR = RGBColor(80, 80, 80)       # Medium gray
TITLE_SIZE = Pt(28)
BODY_SIZE = Pt(36)
CHILD_SIZE = Pt(32)

# --- Parse the markdown into slides ---
# Each top-level bullet (- ) starts a new slide.
# Indented bullets (  - ) are children of the previous top-level bullet.

with open(INPUT_FILE, 'r') as f:
    lines = f.readlines()

title_line = lines[0].strip() if lines else ''

slides_data = []
current_slide = None

for line in lines[1:]:
    stripped = line.rstrip()
    if not stripped:
        continue

    # Top-level bullet: starts with "- " (no leading whitespace)
    if re.match(r'^- ', stripped):
        if current_slide is not None:
            slides_data.append(current_slide)
        current_slide = {
            'body': stripped[2:].strip(),
            'children': []
        }
    # Child bullet: starts with whitespace then "- "
    elif re.match(r'^\s+- ', stripped) and current_slide is not None:
        child_text = re.sub(r'^\s+- ', '', stripped).strip()
        current_slide['children'].append(child_text)

if current_slide is not None:
    slides_data.append(current_slide)

print(f"Parsed {len(slides_data)} announcements from '{INPUT_FILE}'")

# --- Build the presentation ---
prs = Presentation()
# Set widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Title slide
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)
tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = title_line
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TITLE_COLOR
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

# Announcement slides
for i, s in enumerate(slides_data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(6.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Main announcement text
    p = tf.paragraphs[0]
    p.text = s['body']
    p.font.size = BODY_SIZE
    p.font.color.rgb = BODY_COLOR
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = Pt(48)

    # Child bullets
    for child in s['children']:
        p = tf.add_paragraph()
        p.text = f"    \u2022  {child}"
        p.font.size = CHILD_SIZE
        p.font.color.rgb = CHILD_COLOR
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(20)
        p.line_spacing = Pt(44)

    print(f"Slide {i+1}: {s['body'][:60]}...")

prs.save(OUTPUT_FILE)
print(f"\nSaved '{OUTPUT_FILE}' with {len(prs.slides)} slides (1 title + {len(slides_data)} announcements)")
